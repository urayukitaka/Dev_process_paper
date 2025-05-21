from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import pandas as pd
import os

class Report:
    def __init__(self,
                 savefilename:str):
        # save name
        self.savefilename = savefilename
        # initialize
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)  # 16:9
        self.prs.slide_height = Inches(7.5)
        self.blank_slide_layout = self.prs.slide_layouts[6]
    # --------------------------------
    # slide and title
    # --------------------------------
    def _add_title(self, slide, topline: str, subline: str) -> None:
        """Common two‑level title box used by every slide."""
        box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.8))
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True

        run_main = tf.paragraphs[0].add_run()
        run_main.text = topline
        run_main.font.size = Pt(20)
        run_main.font.bold = True
        run_main.font.name = "メイリオ"

        run_sub = tf.add_paragraph().add_run()
        run_sub.text = subline
        run_sub.font.size = Pt(16)
        run_sub.font.name = "メイリオ"

    def _add_picture_if_exists(self, slide, path: str, left, top, *, width=None, height=None) -> None:
        if path and os.path.exists(path):
            slide.shapes.add_picture(path, left, top, width=width, height=height)

    # --------------------------------
    # slide
    # --------------------------------
    def add_analysis_summary(self,
                             topcomment:str,
                             subcomment:str,
                             left_matrix:pd.DataFrame,
                             right_top_graph:str,
                             right_btm_graph:str
                            ) -> None:
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # left matrix
        # --------------------------------
        top_margin = Inches(1.0)  # タイトル下
        # Title
        titlebox = slide.shapes.add_textbox(Inches(0.3), top_margin, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "Correlation Result"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # Matrix
        table_left = Inches(0.3)
        table_top = Inches(1.4)
        table_width_in = self.prs.slide_width.inches / 2 - 0.5
        table_height_in = 5.5
        rows, cols = left_matrix.shape
        table = slide.shapes.add_table(
            rows + 1,
            cols,
            table_left,
            table_top,
            Inches(table_width_in),
            Inches(table_height_in),
        ).table

        # flexible column‑width ratios (fallback to equal widths)
        default_ratios = [0.125, 0.45, 0.15, 0.15]
        if len(default_ratios) != cols:
            default_ratios = [1 / cols] * cols
        total_ratio = sum(default_ratios)
        for idx, ratio in enumerate(default_ratios):
            table.columns[idx].width = Inches((ratio / total_ratio) * table_width_in)

        # color setting
        navy = RGBColor(0, 32, 96)  # 濃い紺色
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)

        # header row
        for c in range(cols):
            cell = table.cell(0, c)
            cell.text = str(left_matrix.columns[c])
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "メイリオ"
            run.font.bold = True
            run.font.size = Pt(8)
            run.font.color.rgb = white
            cell.fill.solid()
            cell.fill.fore_color.rgb = navy

        # data rows
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r + 1, c)
                cell.text = str(left_matrix.iat[r, c])
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = "メイリオ"
                run.font.size = Pt(8)
                run.font.color.rgb = black
                cell.fill.solid()
                cell.fill.fore_color.rgb = white

        # --------------------------------
        # right graphs
        # --------------------------------
        graph_w_in = self.prs.slide_width.inches / 2 - 0.6
        graph_h_in = graph_w_in * 0.4
        graph_w = Inches(graph_w_in)
        graph_h = Inches(graph_h_in)
        right_x = Inches(self.prs.slide_width.inches / 2 + 0.3)
        top_slots = [Inches(1.4), Inches(1.4 + graph_h_in + 0.5)]

        for idx, p in enumerate([right_top_graph, right_btm_graph]):
            # Title
            titlebox = slide.shapes.add_textbox(right_x, top_slots[idx], Inches(3), Inches(0.3))
            if idx==0:
                titlebox.text_frame.text = "Boxplot, SpearmanR by EQP&CHM"
            elif idx==1:
                titlebox.text_frame.text = "Boxplot, SpearmanR by STEP"
            titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
            titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
            self._add_picture_if_exists(slide, p, right_x, top_slots[idx]+Inches(0.3), width=graph_w, height=graph_h)

        self.save_pptx()

    def add_analysis(self,
                     topcomment:str,
                     subcomment:str,
                     top_graph_comment_list:list, # 2 items
                     top_graph_path_list:list, # 2 items
                     btm_graph_comment_list:list, # 22 items
                     btm_graph_path_list:list, # 22 items
                     ):
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # Top graphs
        # --------------------------------
        top_margin = Inches(1.0)  # タイトル下
        graph_width_top = (self.prs.slide_width - Inches(0.6)) / 2
        graph_height_top = graph_width_top * 0.4  # アスペクト比10:4に調整
        left_positions = [Inches(0.3), Inches(0.3) + graph_width_top + Inches(0.1)]
        graph_comments_top = top_graph_comment_list
        graph_paths_top = top_graph_path_list
        for i, left in enumerate(left_positions):
            # Title
            titlebox = slide.shapes.add_textbox(left, top_margin, graph_width_top, Inches(0.3))
            titlebox.text_frame.text = graph_comments_top[i]
            titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
            titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
            # image
            graph_path = graph_paths_top[i]
            if os.path.exists(graph_path):
                slide.shapes.add_picture(graph_path, left, top_margin + Inches(0.2), width=graph_width_top, height=graph_height_top)
        # --------------------------------
        # Bottom graphs
        # --------------------------------
        start_top = top_margin + graph_height_top + Inches(0.4)
        titlebox = slide.shapes.add_textbox(Inches(0.3), start_top, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "CorrelationGraph"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # graphs
        rows, cols = 2, 11
        start_top2 = top_margin + Inches(0.5) + graph_height_top + Inches(0.4)
        # 1:1アスペクト比を優先して左右・上下にマージンを確保
        cell_size = min(
            (self.prs.slide_width - Inches(0.6)) / cols,
            (self.prs.slide_height - start_top2 - Inches(0.5)) / rows
        )
        margin_left = (self.prs.slide_width - (cell_size * cols)) / 2
        margin_top = start_top2
        graph_paths_bottom = btm_graph_path_list
        graph_comments_bottom = btm_graph_comment_list

        for r in range(rows):
            for c in range(cols):
                idx = r * cols + c
                left = margin_left + c * cell_size
                top = margin_top + r * cell_size
                # graph title
                titlebox = slide.shapes.add_textbox(left, top, cell_size, Inches(0.2))
                titlebox.text_frame.text = graph_comments_bottom[idx]
                titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
                # image
                graph_path = graph_paths_bottom[idx]
                if os.path.exists(graph_path):
                    slide.shapes.add_picture(graph_path, left, top + Inches(0.15), width=cell_size, height=cell_size - Inches(0.15))
        # save
        self.save_pptx()

    def save_pptx(self):
        self.prs.save(self.savefilename)

if __name__ == "__main__":
    savefilename = r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\output_slide_fitted.pptx"
    # instance
    REP = Report(savefilename=savefilename)
    # slide 1
    topcomment = "Title, Ranking ABC"
    subcomment = "Sub title, contents"
    df = pd.read_csv(r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\sample_matrix.csv", encoding="CP932")
    df = df.iloc[:22,:] # ranking 22
    right_top_graph=r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\graph_10_6.png"
    right_btm_graph=r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\graph_10_6.png"
    REP.add_analysis_summary(topcomment=topcomment,
                             subcomment=subcomment,
                             left_matrix=df,
                             right_top_graph=right_top_graph,
                             right_btm_graph=right_btm_graph)
    for i in range(10):
        # slide 2
        topcomment = "Title, Report ABC"
        subcomment = "Sub title, contents"
        top_graph_comment_list = ["TrendGraph", "Cum TrendGraph"]
        top_graph_path_list = [r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\graph_10_6.png",
                            r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\graph_10_6.png"]
        btm_graph_comment_list = [f"Equip:{i}" for i in range(1,23)]
        btm_graph_path_list = [r"C:\Users\yktkk\Desktop\DS_practice\programing\pptx\20250521\graph_11_10.png" for r in range(len(btm_graph_comment_list))]
        REP.add_analysis(
            topcomment=topcomment,
            subcomment=subcomment,
            top_graph_comment_list=top_graph_comment_list,
            top_graph_path_list=top_graph_path_list,
            btm_graph_comment_list=btm_graph_comment_list,
            btm_graph_path_list=btm_graph_path_list
        )